"""
LittleJoys Platform Overview — PowerPoint Generator
====================================================
Generates docs/LittleJoys_Platform_Overview.pptx using python-pptx.

Usage:
    python docs/generate_ppt.py

Requirements:
    pip install python-pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ─── Constants ────────────────────────────────────────────────────────────────

OUT_PATH = Path(__file__).parent / "LittleJoys_Platform_Overview.pptx"

# Slide size: Widescreen 13.33" x 7.5"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Brand colours
CORAL    = RGBColor(0xFF, 0x6B, 0x6B)
NAVY     = RGBColor(0x1A, 0x23, 0x40)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF7, 0xF8, 0xFC)
SLATE    = RGBColor(0x2E, 0x3A, 0x52)
MID      = RGBColor(0x4A, 0x55, 0x68)
LIGHT    = RGBColor(0x88, 0x96, 0xA8)
TEAL     = RGBColor(0x4E, 0xCD, 0xC4)
GREEN    = RGBColor(0x2E, 0xCC, 0x71)
TEAL2    = RGBColor(0x02, 0x80, 0x90)
PURPLE   = RGBColor(0x6C, 0x34, 0x83)
DKGREEN  = RGBColor(0x1E, 0x84, 0x49)
ORANGE   = RGBColor(0xE6, 0x7E, 0x22)

FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

# ─── Helper utilities ─────────────────────────────────────────────────────────


def _in(v: float) -> Emu:
    """Convert inches to EMU."""
    return Inches(v)


def add_rect(slide, x, y, w, h, fill: RGBColor, line_color: RGBColor | None = None,
             line_width_pt: float = 0.0) -> None:
    """Add a solid filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _in(x), _in(y), _in(w), _in(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()


def add_textbox(slide, text: str, x, y, w, h, *,
                font_name: str = FONT_BODY,
                font_size: float = 14,
                bold: bool = False,
                italic: bool = False,
                color: RGBColor = SLATE,
                align: PP_ALIGN = PP_ALIGN.LEFT,
                wrap: bool = True,
                ) -> None:
    """Add a simple single-paragraph textbox."""
    txb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_bullets(slide, items: list[str], x, y, w, h, *,
                font_size: float = 13,
                color: RGBColor = SLATE,
                indent: bool = True,
                ) -> None:
    """Add a bulleted list textbox."""
    txb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if indent:
            from pptx.util import Pt as _Pt
            p.level = 0
        run = p.add_run()
        run.text = ("• " if indent else "") + item
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def new_slide(prs: Presentation) -> object:
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def dark_bg(slide, accent_bar: bool = True) -> None:
    """Full dark navy background with optional left coral bar."""
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    if accent_bar:
        add_rect(slide, 0, 0, 0.22, 7.5, CORAL)


def light_bg(slide) -> None:
    """Off-white background with thin coral top bar."""
    add_rect(slide, 0, 0, 13.33, 7.5, OFFWHITE)
    add_rect(slide, 0, 0, 13.33, 0.08, CORAL)


def slide_title(slide, title: str) -> None:
    """Add standard content slide title + divider."""
    add_textbox(slide, title, 0.45, 0.18, 12.4, 0.65,
                font_name=FONT_TITLE, font_size=28, bold=True, color=NAVY,
                align=PP_ALIGN.LEFT)
    add_rect(slide, 0.45, 0.87, 12.4, 0.018, CORAL)


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────


def slide_1_title(prs: Presentation) -> None:
    s = new_slide(prs)
    dark_bg(s, accent_bar=True)

    add_textbox(s, "LittleJoys Persona Simulation Engine", 0.55, 1.35, 9.8, 1.6,
                font_name=FONT_TITLE, font_size=44, bold=True, color=WHITE)
    add_textbox(s, "A Synthetic Research Platform for Kids Nutrition D2C in India",
                0.55, 3.2, 9.8, 0.65,
                font_name=FONT_BODY, font_size=20, color=TEAL)
    add_textbox(s, "Understand 200 Indian parent households — without a single survey",
                0.55, 4.0, 9.8, 0.45,
                font_name=FONT_BODY, font_size=15, italic=True, color=LIGHT)
    add_rect(s, 0.55, 4.5, 5.5, 0.012, CORAL)

    # Badge row
    badges = ["200 Households", "4 Scenarios", "145 Attributes", "Phase A → C"]
    for i, b in enumerate(badges):
        bx = 0.55 + i * 2.05
        add_rect(s, bx, 6.55, 1.85, 0.38, CORAL)
        add_textbox(s, b, bx, 6.55, 1.85, 0.38,
                    font_name=FONT_BODY, font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)


def slide_2_problem(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "The Business Questions We Answer")

    questions = [
        "Which parent segments will actually buy a kids nutrition product?",
        "Where does the purchase funnel break down — and why?",
        "Which marketing or product interventions create the most lift?",
        "How do we test hypotheses on 200 synthetic households in minutes, not months?",
    ]

    for i, q in enumerate(questions):
        # Circle number
        add_rect(s, 0.45, 1.38 + i * 0.92, 0.4, 0.4, CORAL)
        add_textbox(s, str(i + 1), 0.45, 1.38 + i * 0.92, 0.4, 0.4,
                    font_name=FONT_TITLE, font_size=15, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, q, 0.98, 1.35 + i * 0.92, 7.5, 0.5,
                    font_name=FONT_BODY, font_size=14, color=SLATE)

    # Context panel
    add_rect(s, 9.1, 1.0, 3.7, 5.85, NAVY)
    add_textbox(s, "Business Context", 9.3, 1.2, 3.3, 0.38,
                font_name=FONT_TITLE, font_size=15, bold=True, color=CORAL,
                align=PP_ALIGN.CENTER)
    ctx = [
        "LittleJoys: D2C kids nutrition brand in India",
        "4 products, 4 unresolved adoption questions",
        "Traditional research: months + lakhs",
        "This platform: seconds + free",
        "Same seed → reproducible results",
    ]
    add_bullets(s, ctx, 9.3, 1.65, 3.3, 5.0, font_size=12.5, color=OFFWHITE)


def slide_3_architecture(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Five-Phase Research Architecture")

    phases = [
        ("1", "Population\nGeneration", CORAL),
        ("2", "Decision\nSimulation", NAVY),
        ("3", "Probing /\nResearch", TEAL2),
        ("4", "Analysis\nLayer", PURPLE),
        ("5", "Intervention\nSimulation", DKGREEN),
    ]
    for i, (num, label, color) in enumerate(phases):
        x = 0.45 + i * 2.5
        add_rect(s, x, 1.05, 2.2, 1.1, color)
        add_textbox(s, num, x, 1.05, 2.2, 0.5,
                    font_name=FONT_TITLE, font_size=28, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, label, x, 1.58, 2.2, 0.5,
                    font_name=FONT_BODY, font_size=11, color=WHITE,
                    align=PP_ALIGN.CENTER)

    descs = [
        ("Phase 1:", "200 synthetic Indian parent households from validated distributions + Gaussian copula + LLM narratives"),
        ("Phase 2:", "4-layer purchase funnel (Need → Awareness → Consideration → Purchase) per persona × scenario"),
        ("Phase 3:", "Hybrid: LLM persona interviews + counterfactual simulations + statistical attribute probes"),
        ("Phase 4:", "SHAP importance, barrier waterfall, cohort segmentation, trajectory clustering, executive summary"),
        ("Phase 5:", "2×2 intervention quadrant (scope × temporality), lift calculation, ranked leaderboard"),
    ]
    for i, (label, desc) in enumerate(descs):
        add_textbox(s, label, 0.45, 2.45 + i * 0.73, 1.4, 0.34,
                    font_name=FONT_BODY, font_size=13, bold=True, color=NAVY)
        add_textbox(s, desc, 1.88, 2.45 + i * 0.73, 10.87, 0.38,
                    font_name=FONT_BODY, font_size=13, color=SLATE)
        if i < 4:
            add_rect(s, 0.45, 2.79 + i * 0.73, 12.3, 0.008, RGBColor(0xE8, 0xED, 0xF3))

    # Footer
    add_rect(s, 0.45, 6.82, 12.3, 0.42, RGBColor(0xFF, 0xF3, 0xF3))
    add_textbox(s, "Each phase feeds the next. The full pipeline runs in under 30 seconds.",
                0.6, 6.82, 12.1, 0.42,
                font_name=FONT_BODY, font_size=13, italic=True, color=CORAL,
                align=PP_ALIGN.LEFT)


def slide_4_population(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "200 Synthetic Households — Built From Real Distributions")

    stats = [("200", "Synthetic\nHouseholds"), ("145", "Attributes\nPer Persona"),
             ("30",  "Deep Personas\nNarratives"), ("42", "Default Seed\nReproducible"), ("<30s", "Full Pipeline")]
    colors = [CORAL, NAVY, TEAL2, PURPLE, DKGREEN]
    for i, ((stat, label), col) in enumerate(zip(stats, colors)):
        x = 0.45 + i * 2.55
        add_rect(s, x, 1.05, 2.3, 1.15, col)
        add_textbox(s, stat, x, 1.12, 2.3, 0.56,
                    font_name=FONT_TITLE, font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, label, x, 1.68, 2.3, 0.42,
                    font_name=FONT_BODY, font_size=11, color=WHITE,
                    align=PP_ALIGN.CENTER)

    items = [
        "Sampled from validated Indian urban parent demographics (city tier, income, age, education, employment)",
        "City tiers: Tier 1 Metro 45%, Tier 2 City 35%, Tier 3 Emerging 20%",
        "Parent age: 22–45, mean 32; children aged 2–14; 1–5 per household",
        "Income: Tier 1 mean ₹18L/yr, Tier 2 ₹12L/yr, Tier 3 ₹7L/yr (truncated normal per tier)",
        "Psychographic attributes via Gaussian copula with validated inter-attribute correlations",
        "Conditional rule engine: working mothers +0.15 time scarcity, first-time parents +0.12 health anxiety",
        "Tier 3 households +0.10 authority bias; joint families +0.15 elder influence weight",
        "Same seed (default=42) → byte-identical 200 personas every run",
    ]
    add_bullets(s, items, 0.45, 2.42, 12.3, 4.2, font_size=13)


def slide_5_persona_anatomy(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "145 Attributes Per Persona — Across 12 Dimensions")

    dims = [
        ("Demographics",   "city_tier · income_lpa · parent_age · family_structure · socioeconomic_class (A1–C2)"),
        ("Health",         "immunity_concern · growth_concern · nutrition_gap_awareness · medical_authority_trust"),
        ("Psychology",     "health_anxiety · risk_tolerance · information_need · decision_speed · mental_bandwidth"),
        ("Values",         "supplement_belief · indie_brand_openness · best_for_my_child · transparency_importance"),
        ("Cultural",       "dietary_culture · ayurveda_affinity · western_brand_trust · community_orientation"),
        ("Media",          "primary_social_platform · ad_receptivity · digital_payment_comfort · discovery_channel"),
        ("Career",         "employment_status · work_hours · perceived_time_scarcity · cooking_time_available"),
        ("Relationships",  "primary_decision_maker · peer_influence · pediatrician_influence · child_pester_power"),
        ("Education",      "education_level · science_literacy · label_reading_habit · ingredient_awareness"),
        ("Lifestyle",      "clean_label_importance · wellness_trend_follower · parenting_philosophy · meal_planning"),
        ("Daily Routine",  "shopping_platform · budget_consciousness · health_spend_priority · price_reference_point"),
        ("Emotional",      "fear_appeal_responsiveness · aspirational_messaging · testimonial_impact · buyer_remorse"),
    ]

    cols = 2
    row_h = 0.52
    for i, (label, fields) in enumerate(dims):
        col = i % cols
        row = i // cols
        x = 0.45 + col * 6.45
        y = 1.05 + row * row_h

        add_rect(s, x, y + 0.06, 1.55, 0.3, CORAL)
        add_textbox(s, label, x, y + 0.06, 1.55, 0.3,
                    font_name=FONT_BODY, font_size=10.5, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, fields, x + 1.65, y, 4.6, 0.46,
                    font_name=FONT_BODY, font_size=10.5, color=MID)

    add_rect(s, 0.45, 6.92, 12.3, 0.35, RGBColor(0xEB, 0xF5, 0xFB))
    add_textbox(s,
                "All continuous attributes: 0–1 unit intervals. All categoricals: validated enums. "
                "Identity layer is frozen (Pydantic v2, extra='forbid').",
                0.6, 6.92, 12.1, 0.35,
                font_name=FONT_BODY, font_size=11, italic=True, color=NAVY)


def slide_6_narratives(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Deep Personas: First-Person Narratives From Claude")

    items = [
        "30 of 200 personas receive a ~400-word narrative generated by Claude",
        "Covers: household context, child health, shopping behaviour, brand beliefs",
        "K-means cluster-stratified selection for psychographic diversity",
        "Cached by SHA-256 prompt hash — re-runs skip LLM calls entirely",
        "Any narrative persona can be interviewed via the dashboard (20-turn session)",
        "Interview guardrails: prevents AI self-disclosure, stays in character",
    ]
    add_bullets(s, items, 0.45, 1.05, 7.2, 3.3, font_size=13.5)

    # Excerpt card
    add_rect(s, 0.45, 4.55, 7.2, 2.65, RGBColor(0xFF, 0xF8, 0xF8),
             line_color=CORAL, line_width_pt=1.0)
    add_textbox(s, "Example Narrative Excerpt", 0.65, 4.65, 6.8, 0.32,
                font_name=FONT_TITLE, font_size=12, bold=True, color=CORAL)
    excerpt = (
        '"Priya is a 34-year-old working mother in Mumbai\'s Bandra neighbourhood. '
        'She reads ingredient labels habitually, cross-references Amazon reviews with '
        'Instagram reels from paediatrician-backed handles, and distrusts anything '
        'that sounds like hollow marketing. She is not brand-loyal but highly '
        'ingredient-loyal: the first question is always \'what\'s actually in this?\'"'
    )
    add_textbox(s, excerpt, 0.65, 5.02, 6.8, 2.08,
                font_name=FONT_BODY, font_size=12.5, italic=True, color=SLATE)

    # Right steps panel
    add_rect(s, 7.9, 1.05, 4.9, 6.18, NAVY)
    add_textbox(s, "Generation Flow", 8.1, 1.22, 4.5, 0.38,
                font_name=FONT_TITLE, font_size=14, bold=True, color=CORAL,
                align=PP_ALIGN.CENTER)
    steps = [
        ("1", "Tier 2 Selection", "K-means → 30 cluster-stratified"),
        ("2", "Narrative Prompt",  "Claude writes 400-word story"),
        ("3", "Cache Check",       "SHA-256 → skip if cached"),
        ("4", "Interview Session", "PM asks → LLM role-plays persona"),
        ("5", "Guardrails",        "No AI disclosure, 20 turns max"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        add_rect(s, 8.15, 1.75 + i * 0.98, 0.35, 0.35, CORAL)
        add_textbox(s, num, 8.15, 1.75 + i * 0.98, 0.35, 0.35,
                    font_name=FONT_BODY, font_size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title, 8.6, 1.73 + i * 0.98, 3.9, 0.26,
                    font_name=FONT_BODY, font_size=12, bold=True, color=WHITE)
        add_textbox(s, desc, 8.6, 1.98 + i * 0.98, 3.9, 0.3,
                    font_name=FONT_BODY, font_size=10.5, color=LIGHT)


def slide_7_decision_engine(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "4-Layer Purchase Funnel — Per Persona × Per Scenario")

    layer_colors = [CORAL, ORANGE, PURPLE, DKGREEN]
    layers = [
        ("0", "Need Recognition",
         "health_anxiety (×0.20) + nutrition_gap_awareness (×0.25) + child_health_proactivity (×0.20) "
         "+ immunity/growth concerns × age relevance factor"),
        ("1", "Awareness",
         "marketing_budget × channel_match + paediatrician boost (+0.15) + school boost (+0.20) "
         "+ influencer boost (+0.10) + WOM mass"),
        ("2", "Consideration",
         "trust (×0.30) + research_depth (×0.20) + cultural_fit (×0.15) "
         "+ brand_openness (×0.20) + risk_factor (×0.15)"),
        ("3", "Purchase",
         "(value_core × benefit_mix) + emotional_pull − price_barrier − effort_barrier  →  clipped to [0, 1]"),
    ]

    for i, (num, label, desc) in enumerate(layers):
        y = 1.12 + i * 0.82
        add_rect(s, 0.45, y, 0.45, 0.42, layer_colors[i])
        add_textbox(s, num, 0.45, y, 0.45, 0.42,
                    font_name=FONT_TITLE, font_size=16, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, label, 1.0, y, 2.2, 0.42,
                    font_name=FONT_BODY, font_size=14, bold=True, color=NAVY)
        add_textbox(s, desc, 3.3, y, 9.45, 0.42,
                    font_name=FONT_BODY, font_size=12, color=MID)
        if i < 3:
            add_rect(s, 0.45, y + 0.42, 12.3, 0.008, RGBColor(0xE8, 0xED, 0xF3))

    # Rejection reasons
    add_rect(s, 0.45, 4.42, 7.8, 1.9, RGBColor(0xFF, 0xF3, 0xF3),
             line_color=CORAL, line_width_pt=0.8)
    add_textbox(s, "Rejection Reasons (labelled at each exit)", 0.65, 4.52, 7.4, 0.3,
                font_name=FONT_TITLE, font_size=12.5, bold=True, color=CORAL)
    reasons = [
        "Layer 0: age_irrelevant  ·  low_need",
        "Layer 1: low_awareness",
        "Layer 2: dietary_incompatible  ·  insufficient_trust  ·  insufficient_research",
        "Layer 3: price_too_high  ·  effort_too_high  ·  insufficient_trust",
    ]
    add_bullets(s, reasons, 0.65, 4.85, 7.4, 1.38, font_size=12)

    # Calibration box
    add_rect(s, 8.5, 4.42, 4.3, 1.9, NAVY)
    add_textbox(s, "Calibration Target", 8.7, 4.55, 3.9, 0.32,
                font_name=FONT_TITLE, font_size=13, bold=True, color=CORAL,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "12–18%", 8.7, 4.9, 3.9, 0.62,
                font_name=FONT_TITLE, font_size=36, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "First-purchase adoption rate\n(binary search calibration)", 8.7, 5.55, 3.9, 0.62,
                font_name=FONT_BODY, font_size=11.5, color=LIGHT,
                align=PP_ALIGN.CENTER)

    add_rect(s, 0.45, 6.55, 12.3, 0.35, RGBColor(0xEB, 0xF5, 0xFB))
    add_textbox(s, "Each layer threshold is configurable per scenario. Calibrated via binary search to match real-world first-purchase rates.",
                0.6, 6.55, 12.1, 0.35,
                font_name=FONT_BODY, font_size=11, italic=True, color=NAVY)


def slide_8_scenarios(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Four Pre-Built Scenarios, Fully Configurable")

    scens = [
        ("Nutrimix 2–6",       "₹599", "2–6",  "Nutrition powder · Immunity/Growth\nPediatrician + Instagram\nLJ Pass · Temporal 6-month", CORAL),
        ("Nutrimix 7–14",      "₹649", "7–14", "Nutrition powder · Focus/Energy\nSchool partnership focus\nBrand extension · Temporal 12-month", NAVY),
        ("MagBites Gummies",   "₹499", "4–12", "Supplement gummies · Sleep/Calm/Focus\nNew category challenge\nAwareness-first · Static", TEAL2),
        ("ProteinMix",         "₹799", "6–14", "Protein supplement · Muscle/Growth\nHigh effort · Sports club channel\nStatic simulation", PURPLE),
    ]
    for i, (name, price, ages, tags, color) in enumerate(scens):
        x = 0.45 + i * 3.15
        # Card outline
        add_rect(s, x, 1.05, 2.95, 1.6, WHITE, line_color=color, line_width_pt=1.2)
        # Top band
        add_rect(s, x, 1.05, 2.95, 0.4, color)
        add_textbox(s, name, x + 0.12, 1.08, 2.7, 0.34,
                    font_name=FONT_TITLE, font_size=13, bold=True, color=WHITE)
        add_textbox(s, price, x + 0.12, 1.5, 1.2, 0.28,
                    font_name=FONT_BODY, font_size=15, bold=True, color=CORAL)
        add_textbox(s, f"Ages {ages}", x + 1.5, 1.5, 1.3, 0.28,
                    font_name=FONT_BODY, font_size=11, color=MID)
        add_textbox(s, tags, x + 0.12, 1.82, 2.7, 0.75,
                    font_name=FONT_BODY, font_size=10.5, color=MID)

    # Config anatomy
    add_rect(s, 0.45, 2.82, 12.3, 1.7, RGBColor(0xF0, 0xF3, 0xF7),
             line_color=RGBColor(0xD5, 0xDC, 0xE8), line_width_pt=0.8)
    add_textbox(s, "Scenario Config Anatomy", 0.65, 2.92, 11.9, 0.32,
                font_name=FONT_TITLE, font_size=13, bold=True, color=NAVY)

    cfg_items = [
        ("ProductConfig",  "name, category, price_inr, age_range, key_benefits, taste_appeal, clean_label_score, effort_to_acquire"),
        ("MarketingConfig","awareness_budget (0–1), channel_mix (sums 1.0), trust_signals, pediatrician/school/influencer flags"),
        ("LJPassConfig",   "monthly_price: ₹299 · discount: 15% · free_trial: 1 month · retention_boost: +10%"),
        ("Thresholds",     "need_recognition: 0.35 · awareness: 0.30 · consideration: 0.40 · purchase: 0.45"),
    ]
    for i, (label, desc) in enumerate(cfg_items):
        add_textbox(s, label, 0.65 + i * 3.05, 3.28, 2.9, 0.28,
                    font_name=FONT_BODY, font_size=12, bold=True, color=CORAL)
        add_textbox(s, desc, 0.65 + i * 3.05, 3.58, 2.9, 0.85,
                    font_name=FONT_BODY, font_size=10.5, color=MID)

    add_rect(s, 0.45, 4.68, 12.3, 0.55, RGBColor(0xFF, 0xF3, 0xF3),
             line_color=CORAL, line_width_pt=0.8)
    add_textbox(s, "LJ Pass: ₹299/month · 15% discount · 1 free trial month · +10% retention boost · –20% churn reduction. Available: Nutrimix 2–6 and 7–14.",
                0.6, 4.68, 12.1, 0.55,
                font_name=FONT_BODY, font_size=12.5, color=SLATE)

    add_bullets(s, [
        "All scenario parameters are configurable in real-time from the dashboard",
        "Modifications use dot-path notation: 'marketing.pediatrician_endorsement': True, 'product.price_inr': 499.0",
    ], 0.45, 5.42, 12.3, 0.9, font_size=13)


def slide_9_probing(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Hybrid Research: LLM Interviews + Simulations + Statistics")

    probe_types = [
        ("Interview Probe", CORAL, [
            "LLM role-plays persona",
            "Structured open question",
            "15 personas (stratified sample)",
            "k-means clusters responses",
            "→ Themes + representative quotes",
        ]),
        ("Simulation Probe", TEAL2, [
            "Scenario modified via dot-path",
            "Re-run against full population",
            "lift = modified − baseline rate",
            "→ Adoption lift %",
        ]),
        ("Attribute Probe", PURPLE, [
            "Adopters vs Rejectors split",
            "Cohen's d effect size per attribute",
            "Sorted by effect magnitude",
            "→ AttributeSplit list",
        ]),
    ]
    for i, (title, color, lines) in enumerate(probe_types):
        x = 0.45 + i * 4.2
        add_rect(s, x, 1.05, 3.95, 3.1, WHITE, line_color=color, line_width_pt=1.5)
        add_rect(s, x, 1.05, 3.95, 0.52, color)
        add_textbox(s, title, x + 0.18, 1.1, 3.65, 0.38,
                    font_name=FONT_TITLE, font_size=14, bold=True, color=WHITE)
        add_bullets(s, lines, x + 0.18, 1.65, 3.65, 2.4, font_size=12.5)

    # Tree structure
    add_rect(s, 0.45, 4.3, 12.3, 2.0, RGBColor(0xF0, 0xF3, 0xF7),
             line_color=RGBColor(0xD5, 0xDC, 0xE8), line_width_pt=0.8)
    add_textbox(s, "Probing Tree Structure", 0.65, 4.4, 11.9, 0.32,
                font_name=FONT_TITLE, font_size=13, bold=True, color=NAVY)
    add_textbox(s, "ProblemStatement → [ Hypothesis 1 → [ Interview Probe · Simulation Probe · Attribute Probe ], Hypothesis 2 → [...] ] → HypothesisVerdicts → TreeSynthesis",
                0.65, 4.78, 11.9, 0.42,
                font_name=FONT_BODY, font_size=12.5, color=SLATE)
    add_textbox(s, "13 business questions · 4 predefined full trees · 9 lightweight trees · Confidence scoring: Interview 40% + Simulation 40% + Attribute 20%",
                0.65, 5.28, 11.9, 0.82,
                font_name=FONT_BODY, font_size=12, color=MID)

    add_rect(s, 0.45, 6.52, 12.3, 0.35, NAVY)
    add_textbox(s, "HypothesisVerdict + TreeSynthesis produced per run — includes confidence ranking, recommended actions, and synthesis narrative.",
                0.6, 6.52, 12.1, 0.35,
                font_name=FONT_BODY, font_size=11.5, color=WHITE,
                align=PP_ALIGN.CENTER)


def slide_10_analysis(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "From Raw Results to Decision-Quality Insights")

    analyses = [
        ("SHAP Feature Importance", CORAL, [
            "LogisticRegression + LinearExplainer on all 145 attributes",
            "Ranked by mean |SHAP| across population",
            "Example: 'Health Worry Level >0.62 → 2.4× more likely to adopt'",
            "Segment-level splits by city_tier and income_bracket",
        ]),
        ("Barrier Waterfall", NAVY, [
            "36% of rejections at Consideration (primary bottleneck)",
            "29% Need Recognition · 23% Purchase · 13% Awareness",
            "Per-stage top 3 rejection reasons ranked by count",
            "Powers waterfall chart on the Results dashboard page",
        ]),
        ("Cohort Segmentation", TEAL2, [
            "Pre-defined: lapsed users, trust skeptics, time-scarce parents",
            "high-need rejecters, first-time buyers, committed users",
            "Trajectory k-means on 12-month active/inactive time series",
            "Cohort IDs used as target_cohort_id in Phase C interventions",
        ]),
        ("Executive Summary", PURPLE, [
            "Claude LLM: PM-ready narrative (Haiku, temp=0.4)",
            "Fields: headline · trajectory · 3 key_drivers",
            "3 recommendations · 2 risk_factors · mock_mode fallback",
            "Input: monthly counts + clusters + top intervention + funnel",
        ]),
    ]

    for i, (title, color, lines) in enumerate(analyses):
        col = i % 2
        row = i // 2
        x = 0.45 + col * 6.3
        y = 1.05 + row * 2.65

        add_rect(s, x, y, 6.05, 2.45, WHITE, line_color=color, line_width_pt=1.2)
        add_rect(s, x, y, 6.05, 0.44, color)
        add_textbox(s, title, x + 0.18, y + 0.07, 5.65, 0.34,
                    font_name=FONT_TITLE, font_size=14, bold=True, color=WHITE)
        add_bullets(s, lines, x + 0.18, y + 0.5, 5.65, 1.88, font_size=12)

    add_rect(s, 0.45, 6.55, 12.3, 0.35, RGBColor(0xFF, 0xF3, 0xF3),
             line_color=CORAL, line_width_pt=0.8)
    add_textbox(s, "All analysis is reproducible — same population + scenario → same SHAP rankings, same barrier distribution, same executive summary (mock mode).",
                0.6, 6.55, 12.1, 0.35,
                font_name=FONT_BODY, font_size=11, italic=True, color=CORAL)


def slide_11_interventions(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Phase C: Test Interventions Before Spending a Rupee")

    # Axis labels
    add_textbox(s, "NON-TEMPORAL (One-shot changes)", 2.85, 1.05, 3.1, 0.38,
                font_name=FONT_BODY, font_size=11.5, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "TEMPORAL (Sustained over months)", 6.15, 1.05, 3.1, 0.38,
                font_name=FONT_BODY, font_size=11.5, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)

    add_textbox(s, "GENERAL\n(All personas)", 0.45, 1.55, 2.25, 1.95,
                font_name=FONT_BODY, font_size=11.5, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "COHORT-SPECIFIC", 0.45, 3.6, 2.25, 1.95,
                font_name=FONT_BODY, font_size=11.5, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)

    def qcell(title, items, x, y, bg, border):
        add_rect(s, x, y, 3.1, 1.95, bg, line_color=border, line_width_pt=1.0)
        add_textbox(s, title, x + 0.15, y + 0.1, 2.8, 0.36,
                    font_name=FONT_TITLE, font_size=12.5, bold=True, color=NAVY)
        add_bullets(s, items, x + 0.15, y + 0.5, 2.8, 1.38, font_size=10.5)

    qcell("General · Non-Temporal",
          ["15% Returning Customer Discount", "Free Sample of New Flavour", "Pediatrician Campaign (+4pp lift)"],
          2.85, 1.55, RGBColor(0xFF, 0xF8, 0xF8), CORAL)
    qcell("General · Temporal",
          ["LJ Pass ₹199/quarter (+2.5pp)", "Monthly Recipe Subscription", "Brand Ambassador Program"],
          6.15, 1.55, RGBColor(0xF0, 0xFF, 0xF4), DKGREEN)
    qcell("Cohort · Non-Temporal",
          ["Lapsed: ₹100 Coupon (+5.7pp)", "First-time: Starter Kit Discount", "Current: Referral Reward ₹150"],
          2.85, 3.6, RGBColor(0xF0, 0xF8, 0xFF), NAVY)
    qcell("Cohort · Temporal",
          ["Lapsed: Day-22 Reminder+Rotation", "Current: Loyalty Tier Program", "First-time: 90-day Habit Nudge"],
          6.15, 3.6, RGBColor(0xF8, 0xF0, 0xFF), PURPLE)

    # Right panel
    add_rect(s, 9.45, 1.55, 3.4, 5.55, NAVY)
    add_textbox(s, "Lift Formula", 9.65, 1.72, 3.0, 0.34,
                font_name=FONT_TITLE, font_size=13, bold=True, color=CORAL,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "lift_pp = intervention_rate\n       − baseline_rate\n\nlift_pct = lift_pp\n         ÷ baseline × 100",
                9.65, 2.12, 3.0, 1.4,
                font_name="Courier New", font_size=11, color=TEAL,
                align=PP_ALIGN.LEFT)
    add_textbox(s, "48 total interventions\nacross 4 scenarios", 9.65, 3.62, 3.0, 0.55,
                font_name=FONT_BODY, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, "Temporal runs use the full\nevent engine: daily stochastic\nfiring of ad exposures,\ndoctor recommendations,\npack depletion reminders,\nand child reactions.",
                9.65, 4.28, 3.0, 2.6,
                font_name=FONT_BODY, font_size=11, color=LIGHT,
                align=PP_ALIGN.CENTER)

    add_rect(s, 0.45, 6.82, 8.85, 0.35, RGBColor(0xF0, 0xF3, 0xF7),
             line_color=RGBColor(0xD5, 0xDC, 0xE8), line_width_pt=0.8)
    add_textbox(s, "Interventions ranked by lift → leaderboard on the Simulate dashboard page.",
                0.6, 6.82, 8.65, 0.35,
                font_name=FONT_BODY, font_size=11, italic=True, color=NAVY)


def slide_12_outputs(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "What the Platform Produces")

    sections = [
        ("Quantitative Outputs", CORAL, [
            "Adoption rates by scenario, city tier, income bracket, employment",
            "Funnel drop-off counts per stage and rejection reason",
            "SHAP variable importance rankings (145 attributes)",
            "Intervention lift percentages (pp and % vs baseline)",
        ]),
        ("Qualitative Outputs", NAVY, [
            "Themed interview clusters with representative quotes",
            "LLM causal insight statements grounded in SHAP evidence",
            "PM-ready executive summary: headline + trajectory + drivers + recommendations + risks",
            "Hypothesis verdicts with confidence scores",
        ]),
        ("Visualisations", TEAL2, [
            "Funnel waterfall chart (population at each layer + reasons)",
            "Psychographic 2D scatter by adoption outcome",
            "Monthly active rate line chart for temporal scenarios",
            "Intervention lift bar chart ranked by quadrant",
        ]),
        ("Exports & Reproducibility", PURPLE, [
            "PDF research report (exec summary + SHAP + barriers + interviews)",
            "JSON population snapshots (all 200 personas fully serialised)",
            "CSV decision results (one row per persona per scenario)",
            "Same seed → identical results — zero variance across runs",
        ]),
    ]

    for i, (title, color, lines) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = 0.45 + col * 6.3
        y = 1.05 + row * 2.72

        add_rect(s, x, y, 6.05, 2.52, WHITE, line_color=color, line_width_pt=1.2)
        add_rect(s, x, y, 6.05, 0.42, color)
        add_textbox(s, title, x + 0.18, y + 0.07, 5.65, 0.32,
                    font_name=FONT_TITLE, font_size=13.5, bold=True, color=WHITE)
        add_bullets(s, lines, x + 0.18, y + 0.48, 5.65, 1.96, font_size=11.5)

    add_rect(s, 0.45, 6.62, 12.3, 0.6, NAVY)
    add_textbox(s,
                "All outputs feed each other: SHAP findings inform intervention selection; "
                "interview themes confirm quantitative signals; executive summary synthesises everything into one PM-ready document.",
                0.6, 6.62, 12.1, 0.6,
                font_name=FONT_BODY, font_size=12, color=WHITE,
                align=PP_ALIGN.CENTER)


def slide_13_tech(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Technology Stack")

    stack = [
        ("Language",       "Python 3.11"),
        ("Data Models",    "Pydantic v2 — strict frozen identity models, extra='forbid', model validators"),
        ("LLM",            "Anthropic Claude — Sonnet (interviews), Haiku (executive summary), Opus (ReportAgent)"),
        ("ML / Stats",     "scikit-learn (LogisticRegression, KMeans), SHAP (LinearExplainer), NumPy, SciPy"),
        ("Visualisation",  "Plotly (all charts), Streamlit (dashboard + multi-page routing)"),
        ("Simulation",     "Custom event engine — daily stochastic dynamics, WOM propagation, habit formation"),
        ("Testing",        "pytest, 700+ tests, 80%+ coverage target, smoke tests across all 4 scenarios"),
        ("Infrastructure", "uv (package management), structlog (structured logging), Pydantic Settings"),
        ("Caching",        "LLM responses cached by SHA-256 in data/.llm_cache/ — zero re-billing on re-runs"),
        ("Serialisation",  "JSON (population in data/population/), file-based, no database required"),
    ]

    for i, (label, value) in enumerate(stack):
        y = 1.12 + i * 0.53
        add_rect(s, 0.45, y + 0.08, 0.06, 0.3, CORAL)
        add_textbox(s, label, 0.63, y, 2.0, 0.34,
                    font_name=FONT_BODY, font_size=12, bold=True, color=NAVY)
        add_textbox(s, value, 2.73, y, 9.92, 0.38,
                    font_name=FONT_BODY, font_size=12, color=MID)
        if i < len(stack) - 1:
            add_rect(s, 0.45, y + 0.46, 12.3, 0.008, RGBColor(0xE8, 0xED, 0xF3))


def slide_14_deployment(prs: Presentation) -> None:
    s = new_slide(prs)
    light_bg(s)
    slide_title(s, "Deployment and Access")

    items = [
        ("Hosting",      "Render cloud (Python web service, auto-deploy from main branch)"),
        ("Demo Mode",    "Runs without API key — mock LLM, fixture personas, zero cost"),
        ("Population",   "Pre-generated and serialised to data/population/ as JSON"),
        ("LLM Cache",    "SHA-256 keyed responses in data/.llm_cache/ — skips API on re-runs"),
        ("Storage",      "File-based only — no database, no external state, fully portable"),
        ("Local Run",    "streamlit run app/streamlit_app.py"),
        ("Dependencies", "Python 3.11, pyproject.toml + uv.lock, runtime.txt"),
    ]

    for i, (label, value) in enumerate(items):
        add_rect(s, 0.45, 1.08 + i * 0.72, 1.8, 0.46, CORAL)
        add_textbox(s, label, 0.45, 1.08 + i * 0.72, 1.8, 0.46,
                    font_name=FONT_BODY, font_size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, value, 2.4, 1.08 + i * 0.72, 6.1, 0.46,
                    font_name=FONT_BODY, font_size=13, color=SLATE)

    # Right panel
    add_rect(s, 8.75, 1.05, 4.05, 5.95, NAVY)
    add_textbox(s, "Quick Start", 8.95, 1.22, 3.65, 0.36,
                font_name=FONT_TITLE, font_size=15, bold=True, color=CORAL,
                align=PP_ALIGN.CENTER)

    cmds = [
        ("git clone <repo>",                   "Clone the repository"),
        ("uv sync",                            "Install all dependencies"),
        ("cp .env.example .env",               "Configure API key"),
        ("python scripts/generate.py",         "Pre-generate 200 personas"),
        ("streamlit run app/streamlit_app.py", "Launch the dashboard"),
    ]
    for i, (cmd, desc) in enumerate(cmds):
        add_rect(s, 8.95, 1.72 + i * 0.95, 3.65, 0.52, RGBColor(0x0D, 0x1A, 0x30),
                 line_color=TEAL, line_width_pt=0.5)
        add_textbox(s, cmd, 9.08, 1.72 + i * 0.95, 3.38, 0.28,
                    font_name="Courier New", font_size=10.5, color=TEAL)
        add_textbox(s, desc, 9.08, 1.99 + i * 0.95, 3.38, 0.22,
                    font_name=FONT_BODY, font_size=10, color=LIGHT)

    add_rect(s, 0.45, 6.62, 8.15, 0.6, RGBColor(0xF0, 0xF3, 0xF7),
             line_color=RGBColor(0xD5, 0xDC, 0xE8), line_width_pt=0.8)
    add_textbox(s,
                "render.yaml configures cloud deployment. runtime.txt pins Python 3.11. "
                "uv.lock ensures dependency reproducibility across environments.",
                0.6, 6.62, 7.95, 0.6,
                font_name=FONT_BODY, font_size=11.5, color=MID)


# ─── MAIN ─────────────────────────────────────────────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_1_title(prs)
    print("  1/14 — Title")
    slide_2_problem(prs)
    print("  2/14 — The Problem")
    slide_3_architecture(prs)
    print("  3/14 — Architecture")
    slide_4_population(prs)
    print("  4/14 — Population Generation")
    slide_5_persona_anatomy(prs)
    print("  5/14 — Persona Anatomy")
    slide_6_narratives(prs)
    print("  6/14 — LLM Narratives")
    slide_7_decision_engine(prs)
    print("  7/14 — Decision Engine")
    slide_8_scenarios(prs)
    print("  8/14 — Scenarios")
    slide_9_probing(prs)
    print("  9/14 — Probing Pipeline")
    slide_10_analysis(prs)
    print(" 10/14 — Analysis Layer")
    slide_11_interventions(prs)
    print(" 11/14 — Intervention Simulation")
    slide_12_outputs(prs)
    print(" 12/14 — Results & Outputs")
    slide_13_tech(prs)
    print(" 13/14 — Technology Stack")
    slide_14_deployment(prs)
    print(" 14/14 — Deployment")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")


if __name__ == "__main__":
    main()
